#!/usr/bin/env python3
"""Build the SkyDiscover BLIS Router Results deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import os

# ── Colors ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG = RGBColor(0x16, 0x16, 0x2B)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
SUBTLE_GRAY = RGBColor(0x44, 0x44, 0x55)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
CARD_BG = RGBColor(0x22, 0x22, 0x3A)

# Framework colors
FW_COLORS = {
    "openevolve": RGBColor(0x00, 0xB4, 0xA0),
    "evox": RGBColor(0x00, 0x7A, 0xCC),
    "adaevolve": RGBColor(0xFF, 0x8C, 0x42),
    "gepa_native": RGBColor(0x9B, 0x59, 0xB6),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, default_size=13, default_color=LIGHT_GRAY, font_name="Calibri"):
    """lines: list of (text, size, bold, color) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT_TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle_icon(slide, left, top, size, color, label, label_size=9):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(label_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Claude + SkyDiscover Workflow
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, DARK_BG)

# Title
add_text(slide1, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "SkyDiscover: AI-Driven Algorithm Discovery", size=28, bold=True, color=WHITE)
add_text(slide1, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
         "Claude as Experiment Orchestrator — Not the Optimizer", size=15, color=ACCENT_TEAL)

# ── Left side: Workflow diagram ──
add_text(slide1, Inches(0.6), Inches(1.5), Inches(5), Inches(0.4),
         "How It Works", size=18, bold=True, color=WHITE)

# Step boxes - vertical flow
steps = [
    ("1. INITIATE", "Claude reads experiment config,\nsets env vars, creates output dirs", ACCENT_BLUE),
    ("2. LAUNCH", "Claude runs skydiscover-run CLI\n(hands off to search framework)", ACCENT_TEAL),
    ("3. MONITOR", "Claude checks logs every 2 min:\nprogress, errors, best scores", ACCENT_ORANGE),
    ("4. ANALYZE", "Claude runs 4 analysis scripts:\ncompare, plot, effort, diffs", ACCENT_PURPLE),
    ("5. SUMMARIZE", "Claude writes analysis.md with\ntables, findings, recommendations", GREEN),
]

y_start = Inches(2.1)
for i, (title, desc, color) in enumerate(steps):
    y = y_start + Inches(i * 0.95)
    # Color bar on left
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.12), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Card
    add_rounded_rect(slide1, Inches(0.85), y, Inches(5.2), Inches(0.75), CARD_BG)
    add_text(slide1, Inches(1.0), y + Inches(0.05), Inches(1.8), Inches(0.3),
             title, size=12, bold=True, color=color)
    add_text(slide1, Inches(2.7), y + Inches(0.05), Inches(3.2), Inches(0.65),
             desc, size=11, color=LIGHT_GRAY)
    # Arrow between steps
    if i < len(steps) - 1:
        arr_y = y + Inches(0.8)
        add_text(slide1, Inches(1.5), arr_y, Inches(0.5), Inches(0.2),
                 "▼", size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── Right side: Key distinction ──
add_text(slide1, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
         "The Key Distinction", size=18, bold=True, color=WHITE)

# Claude box
add_rounded_rect(slide1, Inches(7.0), Inches(2.1), Inches(5.5), Inches(2.2), CARD_BG, ACCENT_BLUE)
add_text(slide1, Inches(7.3), Inches(2.2), Inches(5), Inches(0.35),
         "Claude (Orchestrator)", size=15, bold=True, color=ACCENT_BLUE)
add_multiline(slide1, Inches(7.3), Inches(2.6), Inches(5), Inches(1.5), [
    ("✓  Sets up experiment environment & config", 12, False, GREEN),
    ("✓  Launches search frameworks via CLI", 12, False, GREEN),
    ("✓  Monitors progress & reports status", 12, False, GREEN),
    ("✓  Runs post-experiment analysis scripts", 12, False, GREEN),
    ("✓  Writes summary with tables & findings", 12, False, GREEN),
    ("✗  Does NOT generate code mutations", 12, False, RED),
    ("✗  Does NOT evaluate candidate programs", 12, False, RED),
])

# Search framework box
add_rounded_rect(slide1, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.2), CARD_BG, ACCENT_TEAL)
add_text(slide1, Inches(7.3), Inches(4.7), Inches(5), Inches(0.35),
         "Search Frameworks (The Actual Optimizers)", size=15, bold=True, color=ACCENT_TEAL)
add_multiline(slide1, Inches(7.3), Inches(5.1), Inches(5), Inches(1.5), [
    ("AdaEvolve — multi-island adaptive evolution", 12, False, ACCENT_ORANGE),
    ("EvoX — evolutionary search with populations", 12, False, ACCENT_BLUE),
    ("OpenEvolve — LLM-guided program synthesis", 12, False, ACCENT_TEAL),
    ("GEPA — genetic programming with archives", 12, False, ACCENT_PURPLE),
    ("", 8),
    ("These frameworks use LLMs to mutate Go code,", 11, False, MID_GRAY),
    ("evaluate via BLIS simulator, and evolve better", 11, False, MID_GRAY),
    ("routing algorithms autonomously.", 11, False, MID_GRAY),
])

# Footer
add_text(slide1, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
         "SkyDiscover Framework  •  BLIS Router Optimization  •  March 2026", size=10, color=MID_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: Experiment Setup
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)

add_text(slide2, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Experiment Setup: BLIS Router Optimization", size=28, bold=True, color=WHITE)
add_text(slide2, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
         "Multi-LLM × Multi-Seed × Multi-Workload evaluation matrix", size=15, color=ACCENT_TEAL)

# ── Top row: 3 cards for multi-* dimensions ──
card_w = Inches(3.7)
card_h = Inches(2.4)
gap = Inches(0.35)
x_start = Inches(0.6)
y_cards = Inches(1.5)

# Card 1: Multi-LLM
add_rounded_rect(slide2, x_start, y_cards, card_w, card_h, CARD_BG, ACCENT_BLUE)
add_text(slide2, x_start + Inches(0.2), y_cards + Inches(0.15), card_w, Inches(0.35),
         "Multi-LLM", size=16, bold=True, color=ACCENT_BLUE)
add_multiline(slide2, x_start + Inches(0.2), y_cards + Inches(0.55), Inches(3.3), Inches(1.7), [
    ("Validates generalization across model sizes", 11, False, LIGHT_GRAY),
    ("", 6),
    ("Model 1: Qwen-7B  (smaller model)", 11, True, WHITE),
    ("  Blackbox estimator, trained coefficients", 10, False, MID_GRAY),
    ("", 4),
    ("Model 2: Qwen-14B  (larger model)", 11, True, WHITE),
    ("  Tests routing under heavier compute", 10, False, MID_GRAY),
    ("", 4),
    ("Scores averaged across both models", 10, False, ACCENT_TEAL),
])

# Card 2: Multi-Seed
x2 = x_start + card_w + gap
add_rounded_rect(slide2, x2, y_cards, card_w, card_h, CARD_BG, ACCENT_ORANGE)
add_text(slide2, x2 + Inches(0.2), y_cards + Inches(0.15), card_w, Inches(0.35),
         "Multi-Seed", size=16, bold=True, color=ACCENT_ORANGE)
add_multiline(slide2, x2 + Inches(0.2), y_cards + Inches(0.55), Inches(3.3), Inches(1.7), [
    ("Prevents overfitting to one random scenario", 11, False, LIGHT_GRAY),
    ("", 6),
    ("Seed 42:  Normal traffic patterns", 11, True, WHITE),
    ("  Standard request arrival distribution", 10, False, MID_GRAY),
    ("", 4),
    ("Seed 456: Bursty traffic stress test", 11, True, WHITE),
    ("  Exposes load-balance saturation", 10, False, MID_GRAY),
    ("", 4),
    ("Scores averaged across both seeds", 10, False, ACCENT_TEAL),
])

# Card 3: Multi-Workload
x3 = x2 + card_w + gap
add_rounded_rect(slide2, x3, y_cards, card_w, card_h, CARD_BG, ACCENT_PURPLE)
add_text(slide2, x3 + Inches(0.2), y_cards + Inches(0.15), card_w, Inches(0.35),
         "Multi-Workload", size=16, bold=True, color=ACCENT_PURPLE)
add_multiline(slide2, x3 + Inches(0.2), y_cards + Inches(0.55), Inches(3.3), Inches(1.7), [
    ("Tests across diverse routing challenges", 11, False, LIGHT_GRAY),
    ("", 6),
    ("cache_warmup: Prefix-affinity vs balance", 11, True, WHITE),
    ("  3 prefix groups across 4 instances", 10, False, MID_GRAY),
    ("", 4),
    ("load_spikes: Bursty arrival routing", 11, True, WHITE),
    ("  One prefix group gets 50% of traffic", 10, False, MID_GRAY),
    ("", 4),
    ("multiturn: Session stickiness", 11, True, WHITE),
    ("  Multi-turn convos with large prefix caches", 10, False, MID_GRAY),
])

# ── Bottom: Single iteration loop diagram ──
add_text(slide2, Inches(0.6), Inches(4.15), Inches(12), Inches(0.4),
         "One Iteration of the Discovery Loop", size=18, bold=True, color=WHITE)

# Flow boxes
loop_y = Inches(4.7)
box_h = Inches(1.8)
box_w = Inches(2.3)
arrow_w = Inches(0.5)

loop_items = [
    ("LLM Generates\nCode Mutation", "Framework asks LLM to\nimprove the Go routing\nfunction (WeightedScoring.\nRoute method)", ACCENT_BLUE),
    ("Go Build\n& Validate", "Compile the mutated\nrouting.go to catch\nsyntax errors before\nrunning simulations", ACCENT_TEAL),
    ("Run 12\nSimulations", "2 seeds × 2 models\n× 3 workloads\n= 12 simulation runs\n(~30-60s total)", ACCENT_ORANGE),
    ("Score &\nArchive", "score = -0.5 × avg_e2e\n  - 0.5 × avg_p95\nKeep if better than\ncurrent best", ACCENT_PURPLE),
    ("Feedback\nto LLM", "Score + error info\nfed back to guide\nnext mutation\n(repeat 50-100×)", GREEN),
]

x_loop = Inches(0.6)
for i, (title, desc, color) in enumerate(loop_items):
    # Box
    add_rounded_rect(slide2, x_loop, loop_y, box_w, box_h, CARD_BG, color)
    # Number circle
    num_shape = slide2.shapes.add_shape(MSO_SHAPE.OVAL, x_loop + Inches(0.08), loop_y + Inches(0.08), Inches(0.35), Inches(0.35))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text(slide2, x_loop + Inches(0.1), loop_y + Inches(0.45), box_w - Inches(0.2), Inches(0.45),
             title, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide2, x_loop + Inches(0.1), loop_y + Inches(0.95), box_w - Inches(0.2), Inches(0.8),
             desc, size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Arrow
    if i < len(loop_items) - 1:
        add_text(slide2, x_loop + box_w, loop_y + Inches(0.7), arrow_w, Inches(0.3),
                 "→", size=22, color=MID_GRAY, align=PP_ALIGN.CENTER)

    x_loop += box_w + arrow_w

# Evaluation matrix callout
add_rounded_rect(slide2, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), CARD_BG)
add_text(slide2, Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.4),
         "Each evaluation = 12 simulations  |  4 frameworks run sequentially  |  50 iterations each  |  ~60s/iteration  |  Total: ~3.5 hours",
         size=12, bold=False, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: Results
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG)

add_text(slide3, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
         "Results: 60-63% Latency Improvement Across All Frameworks", size=28, bold=True, color=WHITE)
add_text(slide3, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
         "Experiment: 260309_50i_twoseed  •  50 iterations  •  Seeds 42,456  •  Multi-LLM ON", size=14, color=MID_GRAY)

# ── Chart: Framework scores (bar chart) ──
from pptx.chart.data import CategoryChartData

chart_data = CategoryChartData()
chart_data.categories = ['Baseline', 'adaevolve', 'openevolve', 'evox', 'gepa_native']
chart_data.add_series('Combined Score', (-14222, -5566, -5265, -5220, -5220))

chart_frame = slide3.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.4), Inches(6.2), Inches(3.5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = False

# Style chart
plot = chart.plots[0]
plot.gap_width = 100
series = plot.series[0]

# Color each bar
bar_colors = [RED, ACCENT_ORANGE, ACCENT_TEAL, ACCENT_BLUE, ACCENT_PURPLE]
for i, color in enumerate(bar_colors):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

# Chart area
chart.chart_style = 2
chart_frame.chart.font.color.rgb = LIGHT_GRAY
chart_frame.chart.font.size = Pt(10)

# Value axis
val_axis = chart.value_axis
val_axis.has_title = True
val_axis.axis_title.text_frame.paragraphs[0].text = "Score (higher = better)"
val_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
val_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = MID_GRAY
val_axis.major_gridlines.format.line.color.rgb = SUBTLE_GRAY

cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(10)

# ── Improvement badges ──
add_text(slide3, Inches(0.6), Inches(5.0), Inches(6.2), Inches(0.3),
         "% Improvement vs Baseline", size=13, bold=True, color=WHITE)

badge_data = [
    ("adaevolve", "+60.9%", ACCENT_ORANGE),
    ("openevolve", "+63.0%", ACCENT_TEAL),
    ("evox", "+63.3%", ACCENT_BLUE),
    ("gepa_native", "+63.3%", ACCENT_PURPLE),
]
bx = Inches(0.6)
for name, pct, color in badge_data:
    add_rounded_rect(slide3, bx, Inches(5.35), Inches(1.45), Inches(0.65), CARD_BG, color)
    add_text(slide3, bx + Inches(0.05), Inches(5.38), Inches(1.35), Inches(0.25),
             name, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide3, bx + Inches(0.05), Inches(5.6), Inches(1.35), Inches(0.3),
             pct, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bx += Inches(1.55)

# ── Right column: Key Findings ──
add_text(slide3, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.35),
         "Key Findings", size=18, bold=True, color=WHITE)

findings = [
    ("All 4 frameworks converge to ~63%", "Similar ceiling suggests a fundamental performance\nboundary for this router configuration.", GREEN),
    ("gepa_native: most sample-efficient", "Best score with only 20 unique programs evaluated\nand lowest build error rate (7% vs 23-30%).", ACCENT_PURPLE),
    ("Multiturn is the differentiator", "Framework scores differ most on multiturn workload\n(325-584ms E2E), while others nearly identical.", ACCENT_ORANGE),
    ("All discover similar core strategies", "Fresh signal priority, non-saturating load penalty,\nKV memory pressure thresholds, session affinity.", ACCENT_BLUE),
]

fy = Inches(1.9)
for title, desc, color in findings:
    bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), fy, Inches(0.08), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide3, Inches(7.45), fy + Inches(0.02), Inches(5.2), Inches(0.25),
             title, size=12, bold=True, color=WHITE)
    add_text(slide3, Inches(7.45), fy + Inches(0.28), Inches(5.2), Inches(0.45),
             desc, size=10, color=LIGHT_GRAY)
    fy += Inches(0.9)

# ── Bottom: Sim-to-Real Transfer ──
add_text(slide3, Inches(7.2), Inches(5.5), Inches(5.5), Inches(0.35),
         "Sim → Real Transfer Viability", size=16, bold=True, color=WHITE)

transfer = [
    ("openevolve", "HIGH", "Deploy as-is. No new signals needed.", GREEN),
    ("gepa_native", "HIGH", "Deploy with SLO middleware.", GREEN),
    ("evox", "MEDIUM", "Needs FreeKVBlocks adapter.", ACCENT_ORANGE),
    ("adaevolve", "LOW-MED", "Hash staleness needs rework.", RED),
]

ty = Inches(5.9)
for fw, verdict, note, color in transfer:
    add_rounded_rect(slide3, Inches(7.2), ty, Inches(5.5), Inches(0.35), CARD_BG)
    add_text(slide3, Inches(7.35), ty + Inches(0.03), Inches(1.3), Inches(0.28),
             fw, size=10, bold=True, color=FW_COLORS.get(fw, WHITE))
    add_text(slide3, Inches(8.6), ty + Inches(0.03), Inches(0.9), Inches(0.28),
             verdict, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide3, Inches(9.5), ty + Inches(0.03), Inches(3.0), Inches(0.28),
             note, size=10, color=LIGHT_GRAY)
    ty += Inches(0.38)

# Footer
add_text(slide3, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
         "Full results: outputs/blis_router/260309_50i_twoseed/analysis.md  •  7 plots + CSV + JSON + diffs available",
         size=10, color=MID_GRAY)

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), "skydiscover_blis_results.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
