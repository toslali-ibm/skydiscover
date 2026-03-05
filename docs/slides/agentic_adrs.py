"""Generate Agentic ADRS slide deck (simple, clean style)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- colors ---
BG = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MUTED = RGBColor(0x88, 0x88, 0x99)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x48, 0x99)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
HEADER_BG = RGBColor(0x0D, 0x47, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def tb(slide, left, top, w, h, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return box


def bullets(slide, left, top, w, h, items, size=14, color=GRAY, spacing=Pt(6)):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return box


def card(slide, left, top, w, h, fill=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def accent_bar(slide, left, top, w, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def table_slide(slide, left, top, rows, col_widths, header_color=HEADER_BG):
    """Add a simple table. rows = list of lists of strings. First row is header."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                        sum(col_widths), Inches(0.35 * n_rows))
    tbl = tbl_shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = GRAY
            # cell fill
            cf = cell.fill
            cf.solid()
            if ri == 0:
                cf.fore_color.rgb = header_color
            else:
                cf.fore_color.rgb = BG_CARD if ri % 2 == 1 else RGBColor(0x1E, 0x28, 0x48)


# =====================================================================
# SLIDE 0A — What is Agentic ADRS? (Version A: side-by-side contrast)
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
   "What is Agentic ADRS?", size=40, color=WHITE, bold=True)
tb(sl, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
   "Both discover better code. The difference is what you learn along the way.",
   size=16, color=MUTED)

# Left card — Vanilla ADRS
card(sl, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
accent_bar(sl, Inches(0.8), Inches(1.8), Inches(5.6), ORANGE)
tb(sl, Inches(1.0), Inches(1.95), Inches(5.2), Inches(0.5),
   "Vanilla ADRS (SkyDiscover / OpenEvolve)", size=18, color=ORANGE, bold=True)
tb(sl, Inches(1.0), Inches(2.5), Inches(5.2), Inches(0.5),
   "LLM generates code variants, evaluator picks the best", size=14, color=GRAY)

bullets(sl, Inches(1.0), Inches(3.1), Inches(5.2), Inches(3.2), [
    "Loop: generate code  >  score it  >  keep winners  >  repeat",
    "LLM is the author — writes mutations and new code",
    "Evaluator is a black-box fitness function",
    "Output: best-scoring code",
    "You learn: what scores highest",
    "You don't learn: why it works",
], size=13, color=GRAY, spacing=Pt(10))

# Right card — Agentic ADRS
card(sl, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8))
accent_bar(sl, Inches(6.8), Inches(1.8), Inches(5.8), ACCENT)
tb(sl, Inches(7.0), Inches(1.95), Inches(5.4), Inches(0.5),
   "Agentic ADRS", size=18, color=ACCENT, bold=True)
tb(sl, Inches(7.0), Inches(2.5), Inches(5.4), Inches(0.5),
   "AI agents form hypotheses, run experiments, extract principles", size=14, color=GRAY)

bullets(sl, Inches(7.0), Inches(3.1), Inches(5.4), Inches(3.2), [
    'Loop: hypothesize  >  design experiment  >  run  >  analyze  >  extract principle',
    "LLM is the scientist — reasons about why, not just what",
    "Evaluator + AI reviewers validate findings",
    "Output: better code + causal principles",
    'You learn: "KV-cache reuse helps because prefix hit rate > 60%"',
    "Principles compound — each iteration builds on past understanding",
], size=13, color=GRAY, spacing=Pt(10))

# Bottom one-liner
tb(sl, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.5),
   "Vanilla = optimization.   Agentic = understanding + optimization.",
   size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 0B — What is Agentic ADRS? (Version B: analogy-driven)
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
   "What is Agentic ADRS?", size=40, color=WHITE, bold=True)

# The analogy
card(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.6))
tb(sl, Inches(1.2), Inches(1.45), Inches(5), Inches(0.5),
   "Vanilla ADRS is like a lab robot", size=20, color=ORANGE, bold=True)
tb(sl, Inches(1.2), Inches(1.95), Inches(5), Inches(0.7),
   "It tries 1000 combinations, tells you which one scored best.\n"
   "You still have to figure out why.", size=14, color=GRAY)

tb(sl, Inches(7.0), Inches(1.45), Inches(5), Inches(0.5),
   "Agentic ADRS is like a research partner", size=20, color=ACCENT, bold=True)
tb(sl, Inches(7.0), Inches(1.95), Inches(5), Inches(0.7),
   'It says: "I think X will work because Y. Let me test it.\n'
   'Here\'s what I found, and what to try next."', size=14, color=GRAY)

# Three key differences
diffs = [
    ("Hypotheses before experiments",
     "Vanilla: mutate code randomly and score it\n"
     "Agentic: form a falsifiable prediction, then design an experiment to test it",
     ACCENT),
    ("AI agents as reviewers, not just generators",
     "Vanilla: LLM writes code, evaluator scores it\n"
     "Agentic: 5 AI reviewers critique the hypothesis, 10 review the findings, self-audit for errors",
     TEAL),
    ("Principles accumulate across iterations",
     "Vanilla: each iteration is independent — no memory of past reasoning\n"
     "Agentic: confirmed principles guide future hypotheses — the system gets smarter over time",
     GREEN),
]

y = Inches(3.2)
for title, desc, clr in diffs:
    card(sl, Inches(0.8), y, Inches(11.7), Inches(1.25))
    accent_bar(sl, Inches(0.8), y, Inches(0.08), clr)
    tb(sl, Inches(1.1), y + Inches(0.08), Inches(4.5), Inches(0.4),
       title, size=15, color=clr, bold=True)
    tb(sl, Inches(1.1), y + Inches(0.45), Inches(11.0), Inches(0.75),
       desc, size=12, color=GRAY)
    y += Inches(1.4)


# =====================================================================
# SLIDE 0C — What is Agentic ADRS? (Version C: minimal / one-pager)
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
   "What is Agentic ADRS?", size=40, color=WHITE, bold=True)
tb(sl, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
   "Use AI agents to run structured experiments on complex systems — not just optimize, but understand.",
   size=16, color=ACCENT)

# Simple 2-column table — that's it
compare_rows = [
    ["", "Vanilla (OpenEvolve / FunSearch)", "Agentic ADRS"],
    ["What the LLM does", "Generates code mutations", "Forms hypotheses, designs experiments"],
    ["What gets evaluated", "Code fitness score", "Whether a prediction was right or wrong"],
    ["What you get back", "Best-scoring code", "Best code + causal principles"],
    ["Memory across iterations", "None — each round is fresh", "Principles compound over time"],
    ["Human role", "Write fitness function, seed code", "Approve hypotheses, validate principles"],
    ["AI role", "Code generator", "Scientist (hypothesize, experiment, reason)"],
    ["Interpretability", "Low — black box search", "High — every change has a causal chain"],
    ["Best for", "Clear objective, need best code fast",
     "Complex systems where understanding matters"],
]
table_slide(sl, Inches(0.8), Inches(1.8), compare_rows,
            [Inches(1.8), Inches(4.8), Inches(5.0)])

# Bottom takeaway
card(sl, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2))
tb(sl, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.4),
   "The key insight", size=16, color=WHITE, bold=True)
tb(sl, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
   'OpenEvolve asks: "what code scores highest?"  Agentic ADRS asks: "why does this mechanism work or fail?"  '
   "One optimizes. The other understands — and then optimizes better.",
   size=14, color=GRAY)


# =====================================================================
# SLIDE 1 — Agentic ADRS vs. Vanilla ADRS
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
   "Agentic ADRS vs. Vanilla ADRS", size=36, color=WHITE, bold=True)

# Left side — what is it + bullets
tb(sl, Inches(0.8), Inches(1.3), Inches(7), Inches(0.5),
   "Hypothesis-driven experimentation with AI agents as reviewers",
   size=16, color=ACCENT)

bullets(sl, Inches(0.8), Inches(1.9), Inches(6.5), Inches(2.2), [
    "A system with many knobs (routing, scheduling, caching, admission)",
    'A hypothesis: "turning knob X will improve metric Y because of reason Z"',
    "Run the experiment",
    "Learn something — whether the hypothesis was right OR wrong",
    "Extract a principle and feed it into the next iteration",
], size=13, color=GRAY)

# Right side — pipeline steps
pipeline = [
    ("Frame the problem", "what is the baseline/target?", ACCENT),
    ("Design hypotheses bundle", "main + ablation + negative\n5 reviewers in parallel", TEAL),
    ("Implement and run", "run.sh & analyze.py\n10 reviews of findings\nself audit", GREEN),
    ("Bayesian tuning", "self audit", ORANGE),
    ("Extract principled & iterate", "confirmed/refuted\nguide future iterations", PINK),
]

x_right = Inches(8.2)
y_step = Inches(1.3)
for label, detail, clr in pipeline:
    card(sl, x_right, y_step, Inches(4.5), Inches(0.85))
    accent_bar(sl, x_right, y_step, Inches(0.08), clr)
    tb(sl, x_right + Inches(0.2), y_step + Inches(0.05), Inches(2.2), Inches(0.35),
       label, size=12, color=clr, bold=True)
    tb(sl, x_right + Inches(2.4), y_step + Inches(0.05), Inches(2.0), Inches(0.75),
       detail, size=10, color=MUTED)
    y_step += Inches(0.95)

# Comparison table
accent_bar(sl, Inches(0.8), Inches(4.35), Inches(11.7), ACCENT)

table_rows = [
    ["", "Strategy Evolution (BLIS)", "OpenEvolve / FunSearch"],
    ["Search method", "Human-guided hypothesis bundles", "LLM-driven evolutionary mutation"],
    ["Unit of change", "A mechanism with causal reasoning", "A code diff scored by fitness"],
    ["Why it works", "Falsifiable predictions + causal models", "Diversity + selection pressure"],
    ["Human role", "Designs hypotheses, approves experiments", "Writes the fitness function, seeds initial code"],
    ["AI role", "Reviews design/code/findings (multi-perspective)", "Generates candidate mutations"],
    ["What you learn", 'Principles ("KV-util is counterproductive because...")', "Better code (but often opaque why)"],
    ["Interpretability", "High — every decision has a causal chain", "Low — evolutionary search is a black box"],
    ["Speed", "Days per iteration", "Hours per iteration"],
    ["Best for", "Understanding complex system interactions", "Optimizing a clear objective function"],
]

table_slide(sl, Inches(0.8), Inches(4.5), table_rows,
            [Inches(1.8), Inches(4.8), Inches(4.8)])

# Footer note
tb(sl, Inches(0.8), Inches(7.05), Inches(11.5), Inches(0.4),
   'The fundamental difference: OpenEvolve asks "what code scores highest?" — '
   'Strategy Evolution asks "why does this mechanism work or fail?" One optimizes, the other understands.',
   size=10, color=MUTED, align=PP_ALIGN.LEFT)


# =====================================================================
# SLIDE 2 — What could be better?
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
   "What could be better?", size=36, color=WHITE, bold=True)

# --- Section 1: Making it structured ---
tb(sl, Inches(0.8), Inches(1.2), Inches(4), Inches(0.4),
   "Making it structured:", size=16, color=ACCENT, bold=True)

struct_rows = [
    ["Gap", "Current", "Improvement"],
    ["Python orchestrator", "Human drives each phase",
     "State machine; each step returns structured JSON; orchestrator decides next action"],
    ["Machine-readable ledger", "Prose/markdown history",
     "JSON/YAML per iteration: {hypothesis, prediction, outcome, effect_size, seeds, metrics, principles}"],
    ["Automated principle checking", "Human reads principles",
     "Executable assertions that gate new designs; auto-check against principle DB"],
    ["Template automation", "Manual file creation",
     "CLI scaffolds bundle directory; 5 arm templates + run.sh + analyze.py stubs"],
    ["Structured scoring", "Pass/fail vs threshold",
     "Multi-objective score (effect size + generalizability + mechanism clarity); Pareto-rank"],
]
table_slide(sl, Inches(0.8), Inches(1.65), struct_rows,
            [Inches(1.8), Inches(2.0), Inches(5.5)])

# --- Section 2: Ideas to borrow ---
tb(sl, Inches(0.8), Inches(4.2), Inches(6), Inches(0.4),
   "Ideas to borrow (from SkyDiscover / OpenEvolve):", size=14, color=TEAL, bold=True)

borrow_rows = [
    ["Idea", "Current", "Adaptation"],
    ["Explore/exploit (UCB)", "Fixed iteration order",
     "Track per-family discovery rate; UCB allocates next iteration to highest-potential family"],
    ["Multi-island populations", "Single linear track",
     "Parallel mechanism families (routing, scheduling, admission, KV); cross-island migration"],
    ["Adaptive population mgmt", "Manual stop criteria",
     "3 consecutive null results = freeze family; contradicted principle = unfocus"],
]
table_slide(sl, Inches(0.8), Inches(4.6), borrow_rows,
            [Inches(2.0), Inches(2.2), Inches(5.5)])

# --- Section 3: Making it generic ---
tb(sl, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
   "Making it generic / replication across codebases", size=16, color=ORANGE, bold=True)

generic_rows = [
    ["Gap", "Current", "Improvement"],
    ["Adaptive Bayesian budget", "Fixed 35-50 evals",
     "Novelty-scaled: similar mechanism = 10 evals warm-started; novel = 50+ evals"],
    ["Statistical effect sizes", ">20% = significant",
     "Cohen's d + confidence intervals for cross-experiment comparison"],
]
table_slide(sl, Inches(0.8), Inches(6.85), generic_rows,
            [Inches(2.0), Inches(2.2), Inches(5.5)])


# =====================================================================
# SLIDE 3 — What kind of paper?
# =====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)

tb(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
   "What kind of paper?", size=36, color=WHITE, bold=True)

# Left side — description
card(sl, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.5))
tb(sl, Inches(1.0), Inches(1.4), Inches(5.4), Inches(0.4),
   "Methodology paper", size=18, color=ACCENT, bold=True)
tb(sl, Inches(1.0), Inches(1.85), Inches(5.4), Inches(0.9),
   "AI-assisted systems experimentation. It's not a new system, not a new algorithm. "
   "It's a new way of using AI agents to systematically explore system design spaces.",
   size=13, color=GRAY)
tb(sl, Inches(1.0), Inches(2.8), Inches(5.4), Inches(0.9),
   'Secondary: experience report — "here\'s what 30 iterations taught us about LLM '
   'inference serving, and the methodology that made those discoveries possible."',
   size=12, color=MUTED)

# Right side — key experiments
card(sl, Inches(7.0), Inches(1.3), Inches(5.8), Inches(2.5))
tb(sl, Inches(7.2), Inches(1.4), Inches(5.4), Inches(0.4),
   "Key experiments", size=16, color=GREEN, bold=True)
bullets(sl, Inches(7.2), Inches(1.85), Inches(5.4), Inches(1.8), [
    "1. Discovery efficiency: How many iterations to find the top-k configurations vs. baselines?",
    "2. Principle quality: Do extracted principles generalize to new workloads / systems?",
    "3. Review gate value: What fraction of CRITICAL findings would have caused wrong conclusions if missed?",
    "4. Prediction accuracy over time: Does the methodology improve prediction accuracy as principles accumulate?",
], size=11, color=GRAY, spacing=Pt(6))

# Comparison baselines
card(sl, Inches(7.0), Inches(4.0), Inches(5.8), Inches(2.0))
tb(sl, Inches(7.2), Inches(4.1), Inches(5.4), Inches(0.4),
   "Comparison baselines", size=16, color=ORANGE, bold=True)
bullets(sl, Inches(7.2), Inches(4.55), Inches(5.4), Inches(1.4), [
    "1. Random search: Same budget (30 iterations), pick random configurations",
    "2. Grid search / parameter sweep: Exhaustive search over the same knob space",
    "3. Bayesian optimization alone: scikit-optimize from scratch, no hypothesis structure",
    "4. OpenEvolve / FunSearch: LLM-driven evolutionary search on the same system",
], size=11, color=GRAY, spacing=Pt(6))

# Venue options table
tb(sl, Inches(0.8), Inches(4.0), Inches(3), Inches(0.4),
   "Venue options", size=16, color=PINK, bold=True)

venue_rows = [
    ["Venue", "Why", "Fit"],
    ["HotOS", "Methodology paper about how to do systems research differently",
     "Good — provocative, short format"],
    ["NSDI / OSDI", "If paired with strong BLIS results showing the methodology found "
     "configurations that baselines missed", "Needs strong empirical results"],
    ["ICSE / FSE (SE venues)", '"AI-assisted software experimentation methodology"',
     "Good fit for the methodology angle"],
    ["NeurIPS / ICML (workshop)", '"LLM agents as research assistants" workshop track',
     "Good for visibility"],
    ["IEEE Software", "Experience report on AI-assisted engineering practices",
     "Good fit, less competitive"],
    ["AAAI (AI+Systems track)", "AI agents doing structured experimentation",
     "Decent fit"],
]
table_slide(sl, Inches(0.8), Inches(4.4), venue_rows,
            [Inches(1.5), Inches(2.8), Inches(1.8)])


# --- save ---
out = "/Users/toslali/Desktop/work/ibm/projects/llm-inference/study/inference-llmd/nsdi/code/skydiscover/docs/slides/agentic_adrs.pptx"
prs.save(out)
print(f"Saved {out}")
print(f"Total slides: {len(prs.slides)}")
