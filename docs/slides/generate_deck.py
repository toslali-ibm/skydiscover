"""Generate SkyDiscover overview slide deck with 2 versions per topic."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- theme colors ----------
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # deep navy
BG_MED = RGBColor(0x16, 0x21, 0x3E)        # slightly lighter navy
ACCENT = RGBColor(0x00, 0xD2, 0xFF)         # sky blue
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)        # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    return shape


def add_tag(slide, left, top, text, color=ACCENT):
    add_rounded_rect(slide, left, top, Inches(1.6), Inches(0.35), color,
                     text=text, font_size=11, font_color=WHITE, bold=True)


def add_version_label(slide, version):
    """Small version label in top-right."""
    add_textbox(slide, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.3),
                f"Version {version}", font_size=10, color=MUTED,
                alignment=PP_ALIGN.RIGHT)


# ======================================================================
# SLIDE 1A - Title slide (version A: bold tagline)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "SkyDiscover", font_size=54, color=ACCENT, bold=True)
add_textbox(sl, Inches(1), Inches(2.8), Inches(10), Inches(0.8),
            "AI-Powered Scientific & Algorithmic Discovery", font_size=28,
            color=WHITE, bold=False)
add_textbox(sl, Inches(1), Inches(4.0), Inches(9), Inches(1.5),
            "A modular framework that uses LLMs to automatically discover,\n"
            "evolve, and optimize code solutions across 200+ benchmarks.",
            font_size=18, color=LIGHT_GRAY)

# tags
tags = ["Evolutionary Search", "LLM-Driven", "200+ Benchmarks", "Pluggable"]
for i, t in enumerate(tags):
    add_tag(sl, Inches(1 + i * 1.8), Inches(5.8), t)

# ======================================================================
# SLIDE 1B - Title slide (version B: question-driven)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(1), Inches(1.2), Inches(11), Inches(1),
            "What if an AI could write better code than you?", font_size=36,
            color=WHITE, bold=True)
add_textbox(sl, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
            "SkyDiscover", font_size=54, color=ACCENT, bold=True)
add_textbox(sl, Inches(1), Inches(4.0), Inches(9), Inches(1.2),
            "Give it a problem + an evaluator.\n"
            "It evolves code solutions using LLMs, keeping what works\n"
            "and improving what doesn't. Fully automated.",
            font_size=18, color=LIGHT_GRAY)

add_bullet_list(sl, Inches(1), Inches(5.5), Inches(10), Inches(1.5),
                ["Math optimization  |  Systems tuning  |  GPU kernels  |  Competitive programming"],
                font_size=16, color=MUTED)

# ======================================================================
# SLIDE 2A - How It Works (version A: step-by-step)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "How SkyDiscover Works", font_size=36, color=WHITE, bold=True)
add_textbox(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "The core loop: Sample \u2192 Prompt \u2192 Generate \u2192 Evaluate \u2192 Store \u2192 Repeat",
            font_size=16, color=ACCENT)

# Five step boxes
steps = [
    ("1. Sample", "Pick a parent\nsolution from\nthe database", ACCENT2),
    ("2. Prompt", "Build an LLM\nprompt with\ncode + feedback", ACCENT),
    ("3. Generate", "LLM writes a\nnew/improved\nversion of the code", GREEN),
    ("4. Evaluate", "Run the evaluator\nto score the\nnew solution", ORANGE),
    ("5. Store", "Save to database\nif it improves\nthe population", PINK),
]
for i, (title, desc, clr) in enumerate(steps):
    x = Inches(0.8 + i * 2.45)
    add_rounded_rect(sl, x, Inches(2.0), Inches(2.2), Inches(2.8), BG_MED,
                     text=f"{title}\n\n{desc}", font_size=15, font_color=WHITE, bold=False)
    # accent bar on top
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()

add_textbox(sl, Inches(0.8), Inches(5.2), Inches(11), Inches(1.8),
            "You provide:  initial_program.py  +  evaluator.py  +  config.yaml\n"
            "SkyDiscover does the rest \u2014 iterating hundreds of times, guided by your evaluator.",
            font_size=15, color=LIGHT_GRAY)

# ======================================================================
# SLIDE 2B - How It Works (version B: inputs/outputs focus)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "How SkyDiscover Works", font_size=36, color=WHITE, bold=True)

# Left: Inputs
add_rounded_rect(sl, Inches(0.8), Inches(1.5), Inches(3.5), Inches(4.5), BG_MED)
add_textbox(sl, Inches(1.0), Inches(1.6), Inches(3), Inches(0.5),
            "YOU PROVIDE", font_size=14, color=ACCENT, bold=True)
items_in = [
    "Initial program\n   (starting code, can be empty)",
    "Evaluator function\n   evaluate(path) \u2192 {score: float}",
    "Config (YAML)\n   algorithm, iterations, LLM model",
]
add_bullet_list(sl, Inches(1.0), Inches(2.2), Inches(3.2), Inches(3.5),
                items_in, font_size=14, color=LIGHT_GRAY, spacing=Pt(14))

# Middle: Arrow
add_textbox(sl, Inches(4.5), Inches(3.3), Inches(1.2), Inches(0.6),
            "\u27A1", font_size=40, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Center: Engine
add_rounded_rect(sl, Inches(5.2), Inches(1.5), Inches(3.5), Inches(4.5), BG_MED)
add_textbox(sl, Inches(5.4), Inches(1.6), Inches(3), Inches(0.5),
            "SKYDISCOVER ENGINE", font_size=14, color=GREEN, bold=True)
engine_items = [
    "LLM generates new code variants",
    "Evaluator scores each variant",
    "Search algorithm picks winners",
    "Repeats 100s\u20131000s of times",
    "Checkpoints progress to disk",
]
add_bullet_list(sl, Inches(5.4), Inches(2.2), Inches(3.2), Inches(3.5),
                engine_items, font_size=14, color=LIGHT_GRAY, spacing=Pt(10))

# Right arrow
add_textbox(sl, Inches(8.9), Inches(3.3), Inches(1.2), Inches(0.6),
            "\u27A1", font_size=40, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Right: Outputs
add_rounded_rect(sl, Inches(9.6), Inches(1.5), Inches(3.2), Inches(4.5), BG_MED)
add_textbox(sl, Inches(9.8), Inches(1.6), Inches(3), Inches(0.5),
            "YOU GET", font_size=14, color=ORANGE, bold=True)
items_out = [
    "Best solution code\n   (best_program.py)",
    "Score + metrics\n   (best_program_info.json)",
    "Full search history\n   (checkpoints, logs)",
]
add_bullet_list(sl, Inches(9.8), Inches(2.2), Inches(3.0), Inches(3.5),
                items_out, font_size=14, color=LIGHT_GRAY, spacing=Pt(14))

# ======================================================================
# SLIDE 3A - Algorithms Overview (version A: table-style)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "Search Algorithms at a Glance", font_size=36, color=WHITE, bold=True)
add_textbox(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "From simple baselines to state-of-the-art co-evolution",
            font_size=16, color=MUTED)

algos = [
    ("TopK", "Keep top K solutions, refine the best", "Baseline", MUTED),
    ("Best-of-N", "Try N variants per round, keep the winner", "Baseline", MUTED),
    ("Beam Search", "Expand a beam with diversity weighting", "Medium", LIGHT_GRAY),
    ("OpenEvolve", "MAP-Elites grid + island migration", "Medium", GREEN),
    ("GEPA", "Pareto-based multi-objective + LLM merge", "Medium", ORANGE),
    ("AdaEvolve", "Multi-island UCB + adaptive intensity + paradigm breakthrough", "Advanced", ACCENT),
    ("EvoX", "Co-evolves solutions AND the search strategy itself", "Advanced", PINK),
]

y = Inches(1.8)
for name, desc, level, clr in algos:
    add_rounded_rect(sl, Inches(0.8), y, Inches(1.8), Inches(0.55), BG_MED,
                     text=name, font_size=15, font_color=clr, bold=True)
    add_textbox(sl, Inches(2.8), y + Inches(0.05), Inches(7), Inches(0.5),
                desc, font_size=14, color=LIGHT_GRAY)
    add_textbox(sl, Inches(10.5), y + Inches(0.05), Inches(2), Inches(0.5),
                level, font_size=12, color=clr, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.7)

# ======================================================================
# SLIDE 3B - Algorithms Overview (version B: grouped cards)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "Search Algorithms at a Glance", font_size=36, color=WHITE, bold=True)

# Three groups
groups = [
    ("Simple", MUTED, [
        ("TopK", "Keep best K, refine one"),
        ("Best-of-N", "N variants per round"),
        ("Beam Search", "Diverse beam expansion"),
    ]),
    ("Quality-Diversity", GREEN, [
        ("OpenEvolve", "MAP-Elites grid\n+ island migration"),
        ("GEPA", "Pareto multi-objective\n+ LLM merge"),
    ]),
    ("Adaptive / Co-Evolving", ACCENT, [
        ("AdaEvolve", "Adaptive intensity\n+ paradigm breakthrough"),
        ("EvoX", "Co-evolves solutions\nAND search strategy"),
    ]),
]

gx = Inches(0.8)
for group_name, gcolor, members in groups:
    gw = Inches(3.8)
    add_textbox(sl, gx, Inches(1.4), gw, Inches(0.4),
                group_name, font_size=16, color=gcolor, bold=True,
                alignment=PP_ALIGN.CENTER)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx, Inches(1.85), gw, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = gcolor
    bar.line.fill.background()

    cy = Inches(2.1)
    for mname, mdesc in members:
        add_rounded_rect(sl, gx + Inches(0.1), cy, gw - Inches(0.2), Inches(1.2), BG_MED)
        add_textbox(sl, gx + Inches(0.3), cy + Inches(0.1), gw - Inches(0.5), Inches(0.4),
                    mname, font_size=15, color=gcolor, bold=True)
        add_textbox(sl, gx + Inches(0.3), cy + Inches(0.5), gw - Inches(0.5), Inches(0.6),
                    mdesc, font_size=13, color=LIGHT_GRAY)
        cy += Inches(1.4)

    gx += Inches(4.1)

# ======================================================================
# SLIDE 4A - AdaEvolve & EvoX Deep Dive (version A: feature list)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "Advanced Algorithms: AdaEvolve & EvoX", font_size=36, color=WHITE, bold=True)

# AdaEvolve card
add_rounded_rect(sl, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_MED)
add_textbox(sl, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
            "AdaEvolve", font_size=24, color=ACCENT, bold=True)
add_textbox(sl, Inches(1.0), Inches(2.1), Inches(5), Inches(0.4),
            "Flagship algorithm \u2014 inspired by biological evolution", font_size=13, color=MUTED)

ada_features = [
    "Multi-Island System \u2014 parallel populations explore different regions",
    "UCB Selection \u2014 smart island picking (explore vs exploit)",
    "Adaptive Intensity \u2014 speeds up when improving, slows when stuck",
    "Paradigm Breakthrough \u2014 LLM proposes entirely new strategies when stuck",
    "Unified Archive \u2014 quality-diversity elite tracking",
    "Error Retry \u2014 feeds errors back to LLM for self-correction",
]
add_bullet_list(sl, Inches(1.0), Inches(2.6), Inches(5.2), Inches(3.8),
                ada_features, font_size=13, color=LIGHT_GRAY, spacing=Pt(10))

# EvoX card
add_rounded_rect(sl, Inches(6.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_MED)
add_textbox(sl, Inches(7.0), Inches(1.5), Inches(5), Inches(0.5),
            "EvoX", font_size=24, color=PINK, bold=True)
add_textbox(sl, Inches(7.0), Inches(2.1), Inches(5), Inches(0.4),
            "Co-evolution \u2014 the search learns to search better", font_size=13, color=MUTED)

evox_features = [
    "Two Populations \u2014 solutions + search strategies evolve together",
    "Strategy as Code \u2014 the search algo is itself a Python program",
    "Strategy Scoring \u2014 strategies rated by how well their solutions improve",
    "Auto Variation Ops \u2014 LLM generates problem-specific mutations",
    "Self-Improving \u2014 better strategies find better solutions, which test strategies",
]
add_bullet_list(sl, Inches(7.0), Inches(2.6), Inches(5.2), Inches(3.8),
                evox_features, font_size=13, color=LIGHT_GRAY, spacing=Pt(10))

# ======================================================================
# SLIDE 4B - AdaEvolve & EvoX Deep Dive (version B: visual mechanism)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "Advanced Algorithms: AdaEvolve & EvoX", font_size=36, color=WHITE, bold=True)

# AdaEvolve mechanism
add_textbox(sl, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
            "AdaEvolve: Adaptive Multi-Island Evolution", font_size=18, color=ACCENT, bold=True)

# Island visualization
island_labels = ["Island 1\n(exploring)", "Island 2\n(exploiting)", "Island 3\n(new paradigm)"]
island_colors = [GREEN, ORANGE, PINK]
for i, (lbl, clr) in enumerate(zip(island_labels, island_colors)):
    x = Inches(0.8 + i * 2.0)
    add_rounded_rect(sl, x, Inches(1.9), Inches(1.8), Inches(1.2), BG_MED,
                     text=lbl, font_size=12, font_color=clr, bold=False)

# arrows between islands
add_textbox(sl, Inches(2.6), Inches(2.2), Inches(0.5), Inches(0.4),
            "\u2194", font_size=20, color=ACCENT)
add_textbox(sl, Inches(4.6), Inches(2.2), Inches(0.5), Inches(0.4),
            "\u2194", font_size=20, color=ACCENT)

add_textbox(sl, Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.0),
            "UCB selects which island to evolve next.\n"
            "Adaptive intensity: fast improvement \u2192 explore more,\n"
            "stagnation \u2192 exploit or trigger paradigm breakthrough.",
            font_size=13, color=LIGHT_GRAY)

# EvoX mechanism
add_textbox(sl, Inches(6.8), Inches(1.3), Inches(6), Inches(0.5),
            "EvoX: Co-Evolution of Solutions & Strategy", font_size=18, color=PINK, bold=True)

# Two linked populations
add_rounded_rect(sl, Inches(6.8), Inches(1.9), Inches(2.5), Inches(1.2), BG_MED,
                 text="Solution\nPopulation", font_size=14, font_color=WHITE, bold=True)
add_rounded_rect(sl, Inches(9.8), Inches(1.9), Inches(2.5), Inches(1.2), BG_MED,
                 text="Strategy\nPopulation", font_size=14, font_color=WHITE, bold=True)

# circular arrows
add_textbox(sl, Inches(9.3), Inches(1.95), Inches(0.6), Inches(0.4),
            "\u21C4", font_size=24, color=PINK)

add_textbox(sl, Inches(6.8), Inches(3.3), Inches(5.5), Inches(1.0),
            "Strategy code decides how to sample & mutate solutions.\n"
            "Solutions test the strategy's effectiveness.\n"
            "Both improve together \u2014 the search learns to search.",
            font_size=13, color=LIGHT_GRAY)

# Comparison box at bottom
add_rounded_rect(sl, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4), BG_MED)
add_textbox(sl, Inches(1.0), Inches(4.7), Inches(11), Inches(0.4),
            "When to use which?", font_size=16, color=WHITE, bold=True)
add_textbox(sl, Inches(1.0), Inches(5.2), Inches(5.3), Inches(1.5),
            "AdaEvolve\n"
            "\u2022 Best for: known problem types\n"
            "\u2022 Strength: robust, well-tuned defaults\n"
            "\u2022 Handles stagnation with paradigm shifts",
            font_size=13, color=LIGHT_GRAY)
add_textbox(sl, Inches(6.8), Inches(5.2), Inches(5.3), Inches(1.5),
            "EvoX\n"
            "\u2022 Best for: novel/unusual problem structures\n"
            "\u2022 Strength: discovers problem-specific strategies\n"
            "\u2022 Higher overhead, higher ceiling",
            font_size=13, color=LIGHT_GRAY)

# ======================================================================
# SLIDE 5A - OpenEvolve & GEPA (version A: side-by-side cards)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "Quality-Diversity Algorithms: OpenEvolve & GEPA", font_size=36, color=WHITE, bold=True)

# OpenEvolve
add_rounded_rect(sl, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_MED)
add_textbox(sl, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
            "OpenEvolve", font_size=24, color=GREEN, bold=True)
add_textbox(sl, Inches(1.0), Inches(2.1), Inches(5), Inches(0.4),
            "Inspired by MAP-Elites \u2014 quality meets diversity", font_size=13, color=MUTED)
oe_features = [
    "MAP-Elites Grid \u2014 cells defined by feature dimensions (complexity, diversity)",
    "Island System \u2014 multiple independent grids with ring migration",
    "Three Modes \u2014 explore (random), exploit (elite), random",
    "Preserves Diversity \u2014 keeps different \"types\" of good solutions",
    "Good for problems with multiple valid approaches",
]
add_bullet_list(sl, Inches(1.0), Inches(2.6), Inches(5.2), Inches(3.8),
                oe_features, font_size=13, color=LIGHT_GRAY, spacing=Pt(10))

# GEPA
add_rounded_rect(sl, Inches(6.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_MED)
add_textbox(sl, Inches(7.0), Inches(1.5), Inches(5), Inches(0.5),
            "GEPA", font_size=24, color=ORANGE, bold=True)
add_textbox(sl, Inches(7.0), Inches(2.1), Inches(5), Inches(0.4),
            "Multi-objective Pareto optimization with LLM merging", font_size=13, color=MUTED)
gepa_features = [
    "Pareto Front \u2014 tracks non-dominated solutions across objectives",
    "Acceptance Gating \u2014 only keeps solutions that improve the front",
    "LLM Merge \u2014 asks LLM to intelligently combine two good solutions",
    "Reflective Prompting \u2014 feeds eval artifacts back into prompts",
    "Best for multi-objective problems (speed vs accuracy, etc.)",
]
add_bullet_list(sl, Inches(7.0), Inches(2.6), Inches(5.2), Inches(3.8),
                gepa_features, font_size=13, color=LIGHT_GRAY, spacing=Pt(10))

# ======================================================================
# SLIDE 5B - OpenEvolve & GEPA (version B: mechanism focus)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "Quality-Diversity Algorithms: OpenEvolve & GEPA", font_size=36, color=WHITE, bold=True)

# OpenEvolve mechanism
add_textbox(sl, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
            "OpenEvolve: MAP-Elites Quality-Diversity", font_size=18, color=GREEN, bold=True)

# Grid visualization
grid_labels = [
    ["Low\ncomplexity\nhigh score", "High\ncomplexity\nhigh score"],
    ["Low\ncomplexity\nlow score", "High\ncomplexity\nlow score"],
]
for r in range(2):
    for c in range(2):
        x = Inches(0.8 + c * 1.6)
        y_pos = Inches(1.9 + r * 1.3)
        clr = GREEN if (r == 0) else MUTED
        add_rounded_rect(sl, x, y_pos, Inches(1.4), Inches(1.1), BG_MED,
                         text=grid_labels[r][c], font_size=10, font_color=clr)

add_textbox(sl, Inches(4.0), Inches(1.9), Inches(2.5), Inches(2.5),
            "Each cell stores the best\nsolution for that \"type\".\n\n"
            "Multiple islands maintain\nindependent grids.\n\n"
            "Top solutions migrate\nbetween islands periodically.",
            font_size=12, color=LIGHT_GRAY)

# GEPA mechanism
add_textbox(sl, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
            "GEPA: Pareto Multi-Objective", font_size=18, color=ORANGE, bold=True)

# Pareto front visualization
add_textbox(sl, Inches(6.8), Inches(1.9), Inches(5.5), Inches(0.5),
            "Objective A (e.g. Speed)  vs  Objective B (e.g. Accuracy)", font_size=12, color=MUTED)

pareto_pts = [
    ("Sol A", Inches(7.0), Inches(2.5)),
    ("Sol B", Inches(8.5), Inches(3.0)),
    ("Sol C", Inches(10.0), Inches(2.8)),
    ("Sol D (merged)", Inches(9.0), Inches(2.3)),
]
for label, px, py in pareto_pts:
    clr = PINK if "merged" in label else ORANGE
    add_rounded_rect(sl, px, py, Inches(1.5), Inches(0.5), BG_MED,
                     text=label, font_size=11, font_color=clr)

add_textbox(sl, Inches(6.8), Inches(3.7), Inches(5.5), Inches(1.5),
            "Only solutions on the Pareto front survive.\n"
            "LLM merges two parents \u2192 potentially better tradeoff.\n"
            "Reflective prompting: eval feedback goes into the next prompt.",
            font_size=12, color=LIGHT_GRAY)

# Bottom comparison
add_rounded_rect(sl, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), BG_MED)
add_textbox(sl, Inches(1.0), Inches(5.3), Inches(11), Inches(0.4),
            "Key Difference", font_size=16, color=WHITE, bold=True)
add_textbox(sl, Inches(1.0), Inches(5.8), Inches(5.3), Inches(1.0),
            "OpenEvolve: Diversity by design \u2014 explicitly maintains\n"
            "different types of solutions in a structured grid.",
            font_size=13, color=GREEN)
add_textbox(sl, Inches(6.8), Inches(5.8), Inches(5.3), Inches(1.0),
            "GEPA: Tradeoff optimization \u2014 finds the best balance\n"
            "across multiple competing objectives.",
            font_size=13, color=ORANGE)

# ======================================================================
# SLIDE 6A - Example Use Case (version A: concrete BLIS example)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "A")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "Example: Optimizing a Load-Balancing Router", font_size=36, color=WHITE, bold=True)
add_textbox(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "Real benchmark: BLIS Router \u2014 routing LLM inference requests across GPU servers",
            font_size=16, color=MUTED)

# Step-by-step
ex_steps = [
    ("1. Define the Problem",
     "You have a Go function that routes requests to servers.\n"
     "Goal: minimize latency while balancing load.",
     ACCENT),
    ("2. Write an Evaluator",
     "evaluator.py runs a simulator, returns:\n"
     "{\"combined_score\": avg_latency_improvement}",
     GREEN),
    ("3. Pick an Algorithm & Run",
     "skydiscover-run routing.go evaluator.py \\\n"
     "  -s adaevolve -i 200 -c config.yaml",
     ORANGE),
    ("4. Get Results",
     "SkyDiscover outputs an optimized routing function\n"
     "that outperforms the hand-tuned baseline.",
     PINK),
]

y = Inches(1.8)
for title, desc, clr in ex_steps:
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.08), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    add_textbox(sl, Inches(1.1), y, Inches(10), Inches(0.4),
                title, font_size=16, color=clr, bold=True)
    add_textbox(sl, Inches(1.1), y + Inches(0.4), Inches(10), Inches(0.7),
                desc, font_size=14, color=LIGHT_GRAY)
    y += Inches(1.35)

# ======================================================================
# SLIDE 6B - Example Use Case (version B: broad use cases)
# ======================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, BG_DARK)
add_version_label(sl, "B")

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
            "What Can You Discover?", font_size=36, color=WHITE, bold=True)
add_textbox(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
            "200+ benchmarks across four domains \u2014 bring your own problem too",
            font_size=16, color=MUTED)

domains = [
    ("Math & Optimization", ACCENT,
     "Function optimization, combinatorics,\nnumerical methods, scheduling"),
    ("Systems & Infrastructure", GREEN,
     "Load balancing, routing algorithms,\ncaching policies, resource allocation"),
    ("GPU Kernels", ORANGE,
     "CUDA/Triton kernel optimization,\nmemory access patterns, tiling strategies"),
    ("Competitive Programming", PINK,
     "Algorithm design, data structures,\ngraph problems, dynamic programming"),
]

for i, (name, clr, desc) in enumerate(domains):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(1.8 + (i // 2) * 2.5)
    add_rounded_rect(sl, x, y, Inches(5.8), Inches(2.0), BG_MED)
    add_textbox(sl, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.5),
                name, font_size=20, color=clr, bold=True)
    add_textbox(sl, x + Inches(0.2), y + Inches(0.7), Inches(5.4), Inches(1.2),
                desc, font_size=15, color=LIGHT_GRAY)

add_textbox(sl, Inches(0.8), Inches(6.9), Inches(11), Inches(0.5),
            "Any problem with code + a scoring function can be a SkyDiscover benchmark.",
            font_size=15, color=ACCENT, alignment=PP_ALIGN.CENTER)


# ---------- save ----------
out_path = "/Users/toslali/Desktop/work/ibm/projects/llm-inference/study/inference-llmd/nsdi/code/skydiscover/docs/slides/skydiscover_overview.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
